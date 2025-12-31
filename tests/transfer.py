import os
import pandas as pd
from docx import Document
from pptx import Presentation
import pdfplumber
from transformers import pipeline


class HFDocumentTranslator:
    def __init__(self):
        print("正在加载中英翻译模型（首次运行会下载 ~300MB，之后离线可用）...")
        # Helsinki-NLP 的中英专用模型，质量高、速度快
        self.translator = pipeline(
            "translation",
            model="Helsinki-NLP/opus-mt-zh-en",
            device=-1  # -1=CPU, 如果有GPU可改为0加速
        )
        print("模型加载完成！")

    def translate_text(self, text):
        if not text or not text.strip():
            return text
        # 批量翻译避免单字太慢
        result = self.translator(text, max_length=512)
        return result[0]['translation_text']

    def translate_excel(self, input_path, output_path=None):
        if output_path is None:
            output_path = input_path.replace('.xlsx', '_en.xlsx')

        df = pd.read_excel(input_path)
        print(f"正在翻译 Excel: {input_path}")

        # 翻译所有文本列
        for col in df.columns:
            if df[col].dtype == 'object':
                print(f"  翻译列: {col}")
                df[col] = df[col].astype(str).apply(self.translate_text)

        df.to_excel(output_path, index=False)
        print(f"翻译完成: {output_path}")
        return output_path

    def translate_word(self, input_path, output_path=None):
        if output_path is None:
            output_path = input_path.replace('.docx', '_en.docx')

        doc = Document(input_path)
        print(f"正在翻译 Word: {input_path}")

        for para in doc.paragraphs:
            if para.text.strip():
                para.text = self.translate_text(para.text)

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        cell.text = self.translate_text(cell.text)

        doc.save(output_path)
        print(f"翻译完成: {output_path}")
        return output_path

    def translate_pdf(self, input_path, output_path=None):
        if output_path is None:
            output_path = input_path.replace('.pdf', '_en.txt')

        print(f"正在翻译 PDF: {input_path}")
        translated_lines = []

        with pdfplumber.open(input_path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    translated = self.translate_text(text)
                    translated_lines.append(f"--- Page {i + 1} ---\n{translated}\n")

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(translated_lines))

        print(f"翻译完成（转为文本）: {output_path}")
        return output_path

    def translate_ppt(self, input_path, output_path=None):
        if output_path is None:
            output_path = input_path.replace('.pptx', '_en.pptx')

        prs = Presentation(input_path)
        print(f"正在翻译 PPT: {input_path}")

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text.strip():
                                run.text = self.translate_text(run.text)

        prs.save(output_path)
        print(f"翻译完成: {output_path}")
        return output_path

    def translate_file(self, file_path):
        ext = os.path.splitext(file_path)[1].lower()
        if ext in ['.xlsx', '.xls']:
            return self.translate_excel(file_path)
        elif ext == '.docx':
            return self.translate_word(file_path)
        elif ext == '.pdf':
            return self.translate_pdf(file_path)
        elif ext == '.pptx':
            return self.translate_ppt(file_path)
        else:
            print(f"不支持格式: {ext}")


# 使用示例
if __name__ == "__main__":
    translator = HFDocumentTranslator()

    # 翻译你的 Excel 文件
    translator.translate_file(r"E:\work\project\YMH-Procyon\Procyon_Software_Functional_Test_Cases - 副本.xlsx")